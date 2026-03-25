"""
File Parser Service
Extrahiert Text aus PDF, DOCX, XLSX, PPTX, TXT, CSV und anderen Formaten.
Alle Parsing-Fehler werden graceful behandelt — nie einen Upload blockieren.
"""
import io
import csv
import chardet
from typing import NamedTuple


class ParseResult(NamedTuple):
    text: str
    page_count: int
    metadata: dict


def detect_encoding(raw_bytes: bytes) -> str:
    """Erkennt die Zeichenkodierung eines byte-Strings."""
    result = chardet.detect(raw_bytes)
    return result.get("encoding") or "utf-8"


def parse_pdf(content: bytes) -> ParseResult:
    """Extrahiert Text aus PDF mit PyMuPDF (fitz)."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(stream=content, filetype="pdf")
        pages = []
        page_count = len(doc)
        for page_num, page in enumerate(doc):
            text = page.get_text("text")
            if text.strip():
                pages.append(f"[Seite {page_num + 1}]\n{text.strip()}")
        doc.close()
        full_text = "\n\n".join(pages)
        return ParseResult(
            text=full_text,
            page_count=page_count,
            metadata={"pdf_pages": page_count}
        )
    except Exception as e:
        return ParseResult(text=f"[PDF-Parsing-Fehler: {e}]", page_count=0, metadata={})


def parse_docx(content: bytes) -> ParseResult:
    """Extrahiert Text aus Word-Dokumenten (.docx)."""
    try:
        from docx import Document
        doc = Document(io.BytesIO(content))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

        # Auch Tabellen auslesen
        table_texts = []
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    table_texts.append(row_text)

        all_text = "\n".join(paragraphs)
        if table_texts:
            all_text += "\n\n[Tabellen]\n" + "\n".join(table_texts)

        return ParseResult(
            text=all_text,
            page_count=1,
            metadata={"paragraphs": len(paragraphs), "tables": len(doc.tables)}
        )
    except Exception as e:
        return ParseResult(text=f"[DOCX-Parsing-Fehler: {e}]", page_count=0, metadata={})


def parse_xlsx(content: bytes) -> ParseResult:
    """Extrahiert Daten aus Excel-Dateien (.xlsx)."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        sheet_texts = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else "" for cell in row]
                # Leere Zeilen überspringen
                if any(cell.strip() for cell in row_data):
                    rows.append(" | ".join(row_data))
            if rows:
                sheet_texts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))

        wb.close()
        return ParseResult(
            text="\n\n".join(sheet_texts),
            page_count=len(wb.sheetnames),
            metadata={"sheets": wb.sheetnames}
        )
    except Exception as e:
        return ParseResult(text=f"[XLSX-Parsing-Fehler: {e}]", page_count=0, metadata={})


def parse_pptx(content: bytes) -> ParseResult:
    """Extrahiert Text aus PowerPoint-Präsentationen (.pptx)."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        slides = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                slides.append(f"[Folie {slide_num}]\n" + "\n".join(slide_texts))

        return ParseResult(
            text="\n\n".join(slides),
            page_count=len(prs.slides),
            metadata={"slides": len(prs.slides)}
        )
    except Exception as e:
        return ParseResult(text=f"[PPTX-Parsing-Fehler: {e}]", page_count=0, metadata={})


def parse_csv(content: bytes) -> ParseResult:
    """Extrahiert Daten aus CSV-Dateien."""
    try:
        encoding = detect_encoding(content)
        text_content = content.decode(encoding, errors="replace")
        reader = csv.reader(io.StringIO(text_content))
        rows = [" | ".join(row) for row in reader if any(cell.strip() for cell in row)]
        return ParseResult(
            text="\n".join(rows),
            page_count=1,
            metadata={"rows": len(rows), "encoding": encoding}
        )
    except Exception as e:
        return ParseResult(text=f"[CSV-Parsing-Fehler: {e}]", page_count=0, metadata={})


def parse_txt(content: bytes) -> ParseResult:
    """Liest Plain-Text-Dateien mit Auto-Encoding-Detection."""
    try:
        encoding = detect_encoding(content)
        text = content.decode(encoding, errors="replace")
        return ParseResult(
            text=text,
            page_count=1,
            metadata={"encoding": encoding, "lines": text.count("\n")}
        )
    except Exception as e:
        return ParseResult(text=f"[TXT-Parsing-Fehler: {e}]", page_count=0, metadata={})


# Mapping: Datei-Endung → Parser-Funktion
EXTENSION_PARSERS = {
    "pdf":  parse_pdf,
    "docx": parse_docx,
    "doc":  parse_docx,    # ältere .doc Dateien (best effort)
    "xlsx": parse_xlsx,
    "xls":  parse_xlsx,    # ältere .xls (best effort via openpyxl)
    "pptx": parse_pptx,
    "ppt":  parse_pptx,
    "csv":  parse_csv,
    "txt":  parse_txt,
    "md":   parse_txt,
    "json": parse_txt,
    "xml":  parse_txt,
    "html": parse_txt,
    "htm":  parse_txt,
    "log":  parse_txt,
}

# MIME-Type Mapping
MIME_MAP = {
    "application/pdf":                                                          "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document":  "docx",
    "application/msword":                                                        "doc",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":        "xlsx",
    "application/vnd.ms-excel":                                                  "xls",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.ms-powerpoint":                                             "ppt",
    "text/csv":                                                                  "csv",
    "text/plain":                                                                "txt",
    "text/markdown":                                                             "md",
    "application/json":                                                          "json",
    "text/xml":                                                                  "xml",
    "application/xml":                                                           "xml",
    "text/html":                                                                 "html",
}

SUPPORTED_EXTENSIONS = set(EXTENSION_PARSERS.keys())
SUPPORTED_MIMES = set(MIME_MAP.keys())


def get_file_extension(filename: str) -> str:
    """Gibt die Datei-Extension in Kleinbuchstaben zurück."""
    return filename.rsplit(".", 1)[-1].lower() if "." in filename else ""


def is_supported(filename: str, mime_type: str = "") -> bool:
    """Prüft ob der Dateityp unterstützt wird."""
    ext = get_file_extension(filename)
    return ext in SUPPORTED_EXTENSIONS or mime_type in SUPPORTED_MIMES


def parse_file(content: bytes, filename: str, mime_type: str = "") -> ParseResult:
    """
    Haupt-Parsing-Funktion.
    Erkennt den Dateityp per Extension (bevorzugt) oder MIME-Type
    und extrahiert den Text.
    """
    ext = get_file_extension(filename)

    # Extension bevorzugen, dann MIME-Type als Fallback
    if ext in EXTENSION_PARSERS:
        parser = EXTENSION_PARSERS[ext]
    elif mime_type in MIME_MAP:
        fallback_ext = MIME_MAP[mime_type]
        parser = EXTENSION_PARSERS.get(fallback_ext, parse_txt)
    else:
        # Unbekannter Typ: als Text versuchen
        parser = parse_txt

    return parser(content)
